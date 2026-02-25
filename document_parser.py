import io
import streamlit as st


def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from a PDF file."""
    try:
        import pdfplumber

        text_parts = []
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
        return "\n\n".join(text_parts)
    except Exception as e:
        st.warning(f"Erro ao ler PDF: {e}")
        return ""


def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from a DOCX file."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(file_bytes))
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        # Also extract from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        return "\n".join(text_parts)
    except Exception as e:
        st.warning(f"Erro ao ler DOCX: {e}")
        return ""


def extract_text_from_txt(file_bytes: bytes) -> str:
    """Extract text from a plain text file."""
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        try:
            return file_bytes.decode("latin-1")
        except Exception as e:
            st.warning(f"Erro ao ler arquivo texto: {e}")
            return ""


def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from a PPTX file."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_bytes))
        text_parts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text.strip())
            if slide_texts:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(text_parts)
    except Exception as e:
        st.warning(f"Erro ao ler PPTX: {e}")
        return ""


def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Extract text from an XLSX file."""
    try:
        import openpyxl

        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    rows.append(row_text)
            if rows:
                text_parts.append(f"[Aba: {sheet_name}]\n" + "\n".join(rows[:100]))  # Limit rows
        return "\n\n".join(text_parts)
    except Exception as e:
        st.warning(f"Erro ao ler XLSX: {e}")
        return ""


PARSERS = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "txt": extract_text_from_txt,
    "md": extract_text_from_txt,
    "pptx": extract_text_from_pptx,
    "xlsx": extract_text_from_xlsx,
}


def parse_uploaded_file(uploaded_file) -> str:
    """Parse an uploaded file and return extracted text."""
    file_bytes = uploaded_file.read()
    extension = uploaded_file.name.rsplit(".", 1)[-1].lower()

    parser = PARSERS.get(extension)
    if parser:
        text = parser(file_bytes)
        if text:
            return f"--- Documento: {uploaded_file.name} ---\n{text}"
        else:
            return f"--- Documento: {uploaded_file.name} (não foi possível extrair texto) ---"
    else:
        st.warning(f"Formato .{extension} não suportado: {uploaded_file.name}")
        return ""


def parse_all_files(uploaded_files) -> str:
    """Parse all uploaded files and return combined text."""
    all_texts = []
    for f in uploaded_files:
        text = parse_uploaded_file(f)
        if text:
            all_texts.append(text)
    return "\n\n".join(all_texts)
